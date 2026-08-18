"""文档解析：把各类文件抽取为带元数据的文本块（page/section）。

返回 blocks: list[{"content": str, "meta": {"page": int|str, "heading": str|None}}]
"""
from __future__ import annotations

import re
from pathlib import Path


def _read_plain(path: Path) -> list[dict]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [{"content": text, "meta": {"page": 1, "heading": None}}]


def _read_markdown(path: Path) -> list[dict]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    blocks: list[dict] = []
    current_heading: str | None = None
    buffer: list[str] = []
    for line in text.splitlines():
        m = re.match(r"^(#{1,6})\s+(.*)$", line)
        if m:
            if buffer:
                blocks.append({"content": "\n".join(buffer),
                               "meta": {"page": len(blocks) + 1, "heading": current_heading}})
                buffer = []
            current_heading = m.group(2).strip()
        else:
            buffer.append(line)
    if buffer:
        blocks.append({"content": "\n".join(buffer),
                       "meta": {"page": len(blocks) + 1, "heading": current_heading}})
    if not blocks:
        blocks = _read_plain(path)
    return blocks


def _read_pdf(path: Path) -> list[dict]:
    blocks: list[dict] = []
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                tables = page.extract_tables()
                for t in tables:
                    if t:
                        text += "\n" + _table_to_markdown(t)
                if text.strip():
                    blocks.append({"content": text.strip(),
                                   "meta": {"page": i, "heading": None}})
    except Exception:
        # 回退 pypdf
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                blocks.append({"content": text.strip(),
                               "meta": {"page": i, "heading": None}})
    return blocks


def _read_docx(path: Path) -> list[dict]:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        rows = [[c.text.strip() for c in row.cells] for row in table.rows]
        parts.append(_table_to_markdown(rows))
    return [{"content": "\n".join(parts), "meta": {"page": 1, "heading": None}}]


def _read_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    blocks: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            blocks.append({"content": "\n".join(texts),
                           "meta": {"page": i, "heading": None}})
    return blocks


def _read_xlsx(path: Path) -> list[dict]:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    blocks: list[dict] = []
    for ws in wb.worksheets:
        rows = [[str(c) if c is not None else "" for c in row] for row in ws.iter_rows(values_only=True)]
        if rows:
            blocks.append({"content": _table_to_markdown(rows),
                           "meta": {"page": ws.title, "heading": None}})
    wb.close()
    return blocks


def _table_to_markdown(rows: list[list]) -> str:
    if not rows:
        return ""
    lines = ["| " + " | ".join(str(c).replace("|", "\\|") for c in rows[0]) + " |"]
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(str(c).replace("|", "\\|") for c in row) + " |")
    return "\n".join(lines)


_SUPPORTED = {
    ".md": _read_markdown, ".markdown": _read_markdown, ".txt": _read_plain,
    ".html": _read_plain, ".htm": _read_plain, ".pdf": _read_pdf,
    ".docx": _read_docx, ".pptx": _read_pptx, ".xlsx": _read_xlsx,
}


def extract_blocks(file_path: str) -> list[dict]:
    path = Path(file_path)
    ext = path.suffix.lower()
    reader = _SUPPORTED.get(ext)
    if reader is None:
        raise ValueError(f"不支持的文件类型: {ext}（支持 {'/'.join(sorted(_SUPPORTED))}）")
    return reader(path)


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in _SUPPORTED
